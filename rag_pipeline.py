import os
import requests
from langchain_community.document_loaders import (
    PyPDFLoader, 
    Docx2txtLoader, 
    TextLoader,
    CSVLoader,
    JSONLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import Chroma
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from dotenv import load_dotenv
import google.generativeai as genai
import logging
from typing import Dict, Any, List, Optional, Union
import json
from pathlib import Path
import mimetypes
import base64
from io import BytesIO

# Additional imports for different file types
try:
    from PIL import Image
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class MultiModalRAGPipeline:
    """Enhanced RAG Pipeline with support for multiple file formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'text': ['.txt', '.md', '.rst'],
        'word': ['.docx', '.doc'],
        'excel': ['.xlsx', '.xls', '.csv'],
        'powerpoint': ['.pptx', '.ppt'],
        'image': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'],
        'json': ['.json'],
        'web': ['.html', '.htm']
    }
    
    def __init__(self, api_key: str = None):
        """Initialize the multi-modal RAG pipeline"""
        self.load_environment(api_key)
        self.setup_embeddings()
        self.setup_llm()
        self.vector_store = None
        self.retrieval_chain = None
        self.check_dependencies()
        
    def check_dependencies(self):
        """Check which optional dependencies are available"""
        deps = {
            'tesseract': TESSERACT_AVAILABLE,
            'docx': DOCX_AVAILABLE,
            'excel': EXCEL_AVAILABLE,
            'powerpoint': PPTX_AVAILABLE
        }
        
        missing = [k for k, v in deps.items() if not v]
        if missing:
            logger.warning(f"Missing optional dependencies: {missing}")
            logger.info("Install with: pip install pillow pytesseract python-docx openpyxl python-pptx pandas")
    
    def load_environment(self, api_key: str = None):
        """Load environment variables and configure API"""
        load_dotenv(dotenv_path=".env.local")
        
        if api_key:
            os.environ["GOOGLE_API_KEY"] = api_key
        
        if "GOOGLE_API_KEY" not in os.environ:
            raise ValueError(
                "GOOGLE_API_KEY not found. Get your API key from "
                "https://aistudio.google.com/app/apikey"
            )
        
        genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
        logger.info("Environment configured successfully")
    
    def setup_embeddings(self):
        """Initialize embedding model"""
        try:
            self.embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            logger.info("Embeddings model initialized")
        except Exception as e:
            logger.error(f"Failed to initialize embeddings: {e}")
            raise
    
    def setup_llm(self, model_name: str = "gemini-1.5-pro"):
        """Initialize LLM"""
        try:
            self.llm = ChatGoogleGenerativeAI(model=model_name)
            logger.info(f"LLM initialized with model: {model_name}")
        except Exception as e:
            logger.error(f"Failed to initialize LLM: {e}")
            raise
    
    def get_file_type(self, file_path: str) -> str:
        """Determine file type based on extension"""
        file_ext = Path(file_path).suffix.lower()
        
        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if file_ext in extensions:
                return file_type
        
        return 'unknown'
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported"""
        return self.get_file_type(file_path) != 'unknown'
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR"""
        if not TESSERACT_AVAILABLE:
            raise ImportError("PIL and pytesseract required for image processing")
        
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from image: {e}")
            return ""
    
    def process_image_with_gemini(self, image_path: str) -> str:
        """Process image using Gemini Vision API"""
        try:
            with open(image_path, 'rb') as image_file:
                image_data = image_file.read()
            
            # Convert to base64 for Gemini
            image_b64 = base64.b64encode(image_data).decode()
            
            # Use Gemini Vision to describe the image
            model = genai.GenerativeModel('gemini-1.5-pro')
            response = model.generate_content([
                "Describe this image in detail, including any text, objects, scenes, and relevant information that would be useful for search and retrieval:",
                {
                    "mime_type": mimetypes.guess_type(image_path)[0] or "image/jpeg",
                    "data": image_b64
                }
            ])
            
            return response.text
            
        except Exception as e:
            logger.error(f"Failed to process image with Gemini: {e}")
            # Fallback to OCR if available
            if TESSERACT_AVAILABLE:
                return self.extract_text_from_image(image_path)
            return ""
    
    def load_pdf_document(self, file_path: str) -> List[Document]:
        """Load PDF document"""
        loader = PyPDFLoader(file_path)
        return loader.load()
    
    def load_text_document(self, file_path: str) -> List[Document]:
        """Load text document"""
        loader = TextLoader(file_path, encoding='utf-8')
        return loader.load()
    
    def load_word_document(self, file_path: str) -> List[Document]:
        """Load Word document"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx required for Word document processing")
        
        loader = Docx2txtLoader(file_path)
        return loader.load()
    
    def load_excel_document(self, file_path: str) -> List[Document]:
        """Load Excel document"""
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl and pandas required for Excel processing")
        
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.csv':
            loader = CSVLoader(file_path)
            return loader.load()
        else:
            # For Excel files, convert to text representation
            try:
                df = pd.read_excel(file_path, sheet_name=None)
                content = ""
                
                for sheet_name, sheet_df in df.items():
                    content += f"\n--- Sheet: {sheet_name} ---\n"
                    content += sheet_df.to_string(index=False)
                    content += "\n"
                
                return [Document(page_content=content, metadata={"source": file_path})]
                
            except Exception as e:
                logger.error(f"Failed to load Excel file: {e}")
                return []
    
    def load_powerpoint_document(self, file_path: str) -> List[Document]:
        """Load PowerPoint document"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required for PowerPoint processing")
        
        try:
            prs = Presentation(file_path)
            content = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"
                
                content += "\n"
            
            return [Document(page_content=content, metadata={"source": file_path})]
            
        except Exception as e:
            logger.error(f"Failed to load PowerPoint file: {e}")
            return []
    
    def load_json_document(self, file_path: str) -> List[Document]:
        """Load JSON document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            content = json.dumps(data, indent=2, ensure_ascii=False)
            return [Document(page_content=content, metadata={"source": file_path})]
            
        except Exception as e:
            logger.error(f"Failed to load JSON file: {e}")
            return []
    
    def load_image_document(self, file_path: str) -> List[Document]:
        """Load and process image document"""
        try:
            # First try Gemini Vision API
            content = self.process_image_with_gemini(file_path)
            
            # If that fails, try OCR
            if not content and TESSERACT_AVAILABLE:
                content = self.extract_text_from_image(file_path)
            
            if content:
                return [Document(
                    page_content=content, 
                    metadata={
                        "source": file_path,
                        "type": "image",
                        "processed_with": "gemini_vision" if "gemini" in content else "ocr"
                    }
                )]
            else:
                logger.warning(f"No text content extracted from image: {file_path}")
                return []
                
        except Exception as e:
            logger.error(f"Failed to process image: {e}")
            return []
    
    def load_document(self, file_path: str) -> List[Document]:
        """Load document based on file type"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_type = self.get_file_type(file_path)
        
        try:
            if file_type == 'pdf':
                docs = self.load_pdf_document(file_path)
            elif file_type == 'text':
                docs = self.load_text_document(file_path)
            elif file_type == 'word':
                docs = self.load_word_document(file_path)
            elif file_type == 'excel':
                docs = self.load_excel_document(file_path)
            elif file_type == 'powerpoint':
                docs = self.load_powerpoint_document(file_path)
            elif file_type == 'json':
                docs = self.load_json_document(file_path)
            elif file_type == 'image':
                docs = self.load_image_document(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            logger.info(f"Loaded {len(docs)} documents from {file_path} (type: {file_type})")
            return docs
            
        except Exception as e:
            logger.error(f"Failed to load document {file_path}: {e}")
            raise
    
    def load_multiple_documents(self, file_paths: List[str]) -> List[Document]:
        """Load multiple documents of different types"""
        all_docs = []
        
        for file_path in file_paths:
            try:
                docs = self.load_document(file_path)
                all_docs.extend(docs)
            except Exception as e:
                logger.error(f"Failed to load {file_path}: {e}")
                continue
        
        logger.info(f"Loaded total {len(all_docs)} documents from {len(file_paths)} files")
        return all_docs
    
    def download_file(self, url: str, output_path: str) -> str:
        """Download file from URL"""
        try:
            if os.path.exists(output_path):
                logger.info(f"{output_path} already exists. Skipping download.")
                return output_path
            
            logger.info(f"Downloading {url}...")
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            with open(output_path, "wb") as f:
                f.write(response.content)
            
            logger.info(f"Download complete: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Failed to download file: {e}")
            raise
    
    def split_documents(self, docs: List[Document], chunk_size: int = 1000, 
                       chunk_overlap: int = 200) -> List[Document]:
        """Split documents into chunks"""
        try:
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, 
                chunk_overlap=chunk_overlap
            )
            documents = text_splitter.split_documents(docs)
            logger.info(f"Split into {len(documents)} chunks")
            return documents
        except Exception as e:
            logger.error(f"Failed to split documents: {e}")
            raise
    
    def create_vector_store(self, documents: List[Document], persist_directory: str = None):
        """Create and store vector embeddings"""
        try:
            if persist_directory:
                self.vector_store = Chroma.from_documents(
                    documents, 
                    self.embeddings,
                    persist_directory=persist_directory
                )
            else:
                self.vector_store = Chroma.from_documents(documents, self.embeddings)
            
            logger.info("Vector store created successfully")
        except Exception as e:
            logger.error(f"Failed to create vector store: {e}")
            raise
    
    def setup_retrieval_chain(self):
        """Set up the retrieval chain"""
        try:
            if not self.vector_store:
                raise ValueError("Vector store not initialized")
            
            # Enhanced prompt for multi-modal content
            question_answering_prompt = ChatPromptTemplate.from_messages([
                ("system", 
                 "You are an assistant for question-answering tasks with access to multi-modal content "
                 "including text documents, images, presentations, and structured data. "
                 "Use the following retrieved context to answer the question. "
                 "If the context includes image descriptions or structured data, incorporate that information appropriately. "
                 "If you don't know the answer, just say that you don't know. "
                 "Provide comprehensive answers but keep them concise and relevant.\n\n"
                 "Context: {context}"),
                ("user", "{input}"),
            ])
            
            # Create document chain
            document_chain = create_stuff_documents_chain(self.llm, question_answering_prompt)
            
            # Create retriever with enhanced search
            retriever = self.vector_store.as_retriever(
                search_type="similarity",
                search_kwargs={"k": 5}  # Return top 5 similar documents
            )
            
            # Create retrieval chain
            self.retrieval_chain = create_retrieval_chain(retriever, document_chain)
            logger.info("Retrieval chain setup complete")
            
        except Exception as e:
            logger.error(f"Failed to setup retrieval chain: {e}")
            raise
    
    def process_documents(self, file_paths: Union[str, List[str]], 
                         chunk_size: int = 1000, chunk_overlap: int = 200, 
                         persist_directory: str = None) -> Dict[str, Any]:
        """Complete document processing pipeline for multiple file types"""
        try:
            # Handle single file or list of files
            if isinstance(file_paths, str):
                file_paths = [file_paths]
            
            # Load documents
            docs = self.load_multiple_documents(file_paths)
            
            if not docs:
                raise ValueError("No documents could be loaded")
            
            # Split into chunks
            documents = self.split_documents(docs, chunk_size, chunk_overlap)
            
            # Create vector store
            self.create_vector_store(documents, persist_directory)
            
            # Setup retrieval chain
            self.setup_retrieval_chain()
            
            # Analyze file types
            file_types = {}
            for file_path in file_paths:
                file_type = self.get_file_type(file_path)
                file_types[file_type] = file_types.get(file_type, 0) + 1
            
            return {
                "num_files": len(file_paths),
                "num_documents": len(docs),
                "num_chunks": len(documents),
                "file_types": file_types,
                "supported_formats": list(self.SUPPORTED_FORMATS.keys()),
                "status": "success"
            }
            
        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            raise
    
    def query(self, question: str) -> Dict[str, Any]:
        """Query the multi-modal RAG system"""
        try:
            if not self.retrieval_chain:
                raise ValueError("Retrieval chain not initialized. Process documents first.")
            
            logger.info(f"Querying: '{question}'")
            response = self.retrieval_chain.invoke({"input": question})
            
            # Enhanced response formatting
            result = {
                "question": question,
                "answer": response.get("answer", "No answer found"),
                "context": [doc.page_content for doc in response.get("context", [])],
                "metadata": [doc.metadata for doc in response.get("context", [])],
                "source_types": list(set(doc.metadata.get("type", "text") for doc in response.get("context", [])))
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Query failed: {e}")
            raise
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get list of supported file formats"""
        return self.SUPPORTED_FORMATS.copy()
    
    def batch_query(self, questions: List[str]) -> List[Dict[str, Any]]:
        """Process multiple questions"""
        results = []
        for question in questions:
            try:
                result = self.query(question)
                results.append(result)
            except Exception as e:
                results.append({
                    "question": question,
                    "answer": f"Error: {str(e)}",
                    "context": [],
                    "metadata": [],
                    "source_types": []
                })
        return results


# Example usage
def main():
    """Example usage of the multi-modal RAG pipeline"""
    try:
        # Initialize pipeline
        rag = MultiModalRAGPipeline()
        
        # Show supported formats
        logger.info("Supported formats:")
        for file_type, extensions in rag.get_supported_formats().items():
            logger.info(f"  {file_type}: {extensions}")
        
        # Example with multiple file types
        sample_files = [
            "sample_document.pdf",
            "sample_image.jpg",
            "sample_document.docx",
            "sample_data.xlsx"
        ]
        
        # Filter to existing files
        existing_files = [f for f in sample_files if os.path.exists(f)]
        
        if existing_files:
            # Process documents
            result = rag.process_documents(existing_files, persist_directory="./multi_modal_db")
            logger.info(f"Processing result: {result}")
            
            # Query examples
            questions = [
                "What is the main content of the documents?",
                "Are there any images? What do they show?",
                "What data is available in the spreadsheets?",
                "Summarize the key information from all sources."
            ]
            
            # Process queries
            for question in questions:
                try:
                    response = rag.query(question)
                    logger.info(f"Q: {question}")
                    logger.info(f"A: {response['answer']}")
                    logger.info(f"Sources: {response['source_types']}")
                    logger.info("-" * 50)
                except Exception as e:
                    logger.error(f"Query failed: {e}")
        else:
            logger.warning("No sample files found for demonstration")
        
    except Exception as e:
        logger.error(f"Main execution failed: {e}")


if __name__ == "__main__":
    main()